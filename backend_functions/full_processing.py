from google.cloud import language_v1
from google.cloud import translate_v2 as translate
import requests
import docx2txt
from pptx import Presentation
from firebase_admin import firestore

class FullProcessor:

    '''
    Full processing of company description and files. 
    
    1) Read company from Firestore DB
    2) Parse Files (PDFs, PPTs, DOCs)
    3) Call Google NLP API
    4) Generate Keywords
    5) Write keywords to Firestore DB

    '''

    def call_ocr(self, file_uri, dest_uri):
        # Single PDF OCR

        #TODO: Insert Google OCR script here with 50 page max limit

        return "hi"

    def parse_files(self, database, fs, filename, file_uri):

        pdfbox_endpoints = {
            'main':'https://us-central1-coopsightsoftware.cloudfunctions.net/backend-pdfbox',
            'staging':'https://us-central1-coopsightstaging.cloudfunctions.net/backend-pdfbox-staging',
            'vc': 'https://us-central1-coopsight-vc.cloudfunctions.net/backend-pdfbox-vc',
            'general': 'https://us-central1-coopsight-general.cloudfunctions.net/backend-pdfbox-general',
            'open': 'https://europe-west3-coopsight-open-network.cloudfunctions.net/backend-pdfbox-open'
        }

        if filename.lower().endswith(('.doc', '.docx')):
            temp_uri = file_uri.split("gs://",1)[1]
            read = fs.open(temp_uri)
            clean_text = docx2txt.process(read)

            return clean_text, True
        
        elif filename.lower().endswith(('.ppt', '.pptx')):
            temp_uri = file_uri.split("gs://",1)[1]
            read = fs.open(temp_uri)
            prs = Presentation(read)
            slide_number = 1
            clean_text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        clean_text += shape.text + "\n"
                slide_number += 1

            return clean_text, True

        elif filename.lower().endswith('.pdf'):
            # print(filename)
            endpoint = pdfbox_endpoints[database]
            temp_uri = file_uri.split("appspot.com/",1)[1]
            pdfbox_params = {'uri':temp_uri}
            pdfbox_response = requests.post(endpoint, json=pdfbox_params)
            # print(pdfbox_response.text)

            # TODO: ADD OCR Integration code here. Call call_ocr function above. 

            return pdfbox_response.text, True

        else:
            return '', False

    def call_nlp(self, nlp_client, test_string, desc_string, web_string):

        # Call desc_string first
        test_string = test_string.replace('\n',' ')
        word_dict = {}
        type_ = language_v1.Document.Type.PLAIN_TEXT
        language = "en"
        encoding_type = language_v1.EncodingType.UTF8

        #* Call NLP API on company description
        document = {"content": desc_string, "type_": type_, "language": language}
        nlp_response = nlp_client.analyze_entities(request = {'document': document, 'encoding_type': encoding_type})
        for entity in nlp_response.entities:
                word_dict[entity.name] = entity.salience

        #* Call NLP API on web string
        if len(web_string) > 150:
            document = {"content": web_string, "type_": type_, "language": language}
            nlp_response = nlp_client.analyze_entities(request = {'document': document, 'encoding_type': encoding_type})
            for entity in nlp_response.entities:
                    word_dict[entity.name] = entity.salience

        #* Set string upper length for test_string
        print(len(test_string))
        divisor = 100000
        count = 0
        for i in range(int(len(test_string)/divisor) + 1):
            # print("call " + str(i))
            temp_string = test_string[i*divisor:i*divisor+divisor]
            # print(len(temp_string))

            document = {"content": temp_string, "type_": type_, "language": language}

            nlp_response = nlp_client.analyze_entities(request = {'document': document, 'encoding_type': encoding_type})

            for entity in nlp_response.entities:
                word_dict[entity.name] = entity.salience

            count += 1
            if count > 1:
                break

        return word_dict


    def generate_keywords(self, nlp_response):

        # word_dict = {}
        # for entity in nlp_response.entities:
        #     word_dict[entity.name] = entity.salience
            
        sorted_words = sorted(nlp_response.items(), key=lambda x: x[1], reverse=True)

        # top 20 and 100
        topk = {}
        topn = {}
        counter = 0
        for w in sorted_words:
            if w[1] > 0:
                if counter < 20:
                    topk[w[0]] = w[1]
                if counter < 100:
                    topn[w[0]] = w[1]
                counter += 1

        return topk, topn

    def translate(self, translate_client, website_string):
        result = translate_client.detect_language(website_string[0:2000])
        if result['language'] == 'en':
            website_string = str(website_string.encode('ascii', errors='ignore').decode('utf-8'))
            return website_string
        else:
            print(f'translating website from {result["language"]}')
            result = translate_client.translate(website_string[0:2000], target_language='en')
            print(result["translatedText"])
            website_string = str(result["translatedText"].encode('ascii', errors='ignore').decode('utf-8'))
            return website_string

    def execute(self, database, firebase_db, nlp_client, translate_client, fs, docid, path, label=None, parent=None, pdfs=None, industry=None, website=None, WS=None):

        #* FETCH COMPANY DESCRIPTION TEXT AT PATH
        #!!! THIS IS CURCIAL <CONTAINTS DEPRECATED CONTENT>
        if path[0] == '/':
            path = path[1:]
        desc_ref = firebase_db.document(path)
        desc_doc = desc_ref.get()
            
        if desc_doc.exists:
            desc_string = desc_doc.to_dict()['description']
        else:
            desc_string = ''

        #* RETRIEVE LIST OF FILES
        #!!! DIDFF IMPLEMENTATION FOR DIFF DATABASE
        if database == 'staging' or database == 'open':
            files_ref = firebase_db.collection('privateBins')
            query_ref = files_ref.where('docid', '==', docid).stream()
        else:
            files_ref = firebase_db.collection('files')
            query_ref = files_ref.where('docid', '==', docid).stream()
        
        count = 0
        for doc in query_ref:
            query_doc = doc.to_dict()
            count += 1
            break

        #* PARSE LIST OF FILES
        parsed_list = []
        if count > 0 and pdfs is None:
            #!!!! THIS IS CURCIAL <CONTAINTS DEPRECATED CONTENT>
            if database == 'general':
                for item in query_doc['filesList']:
                    word_string, status = self.parse_files(database, fs, item['fileName'], item['fileURI'])

                    if status == True:
                        parsed_list.append(word_string)
                    else:
                        print("failed to parse " + item['fileName'])
            elif database == 'staging' or database == 'open':
                for item in query_doc['filesList']:
                    word_string, status = self.parse_files(database, fs, item['name'], item['fileURI'])
                    if status == True:
                        parsed_list.append(word_string)
                    else:
                        print("failed to parse " + item['name'])
        
        #* PARSE LIST OF FILES FROM INPUT PARAM
        if pdfs is not None:
            for pdf in pdfs:
                word_string, status = self.parse_files(database, fs, '.pdf', pdf)
                if status == True:
                    parsed_list.append(word_string)
                else:
                    print("failed to parse " + '.pdf')

        #* INSERT STRINGS FROM WEB SCRAPING
        web_string = ''
        if website is not None and WS is not None:
            web_string = WS.scrape(website, False)

        #* TEST DETECTION OF LANGUAGE
        web_string = self.translate(translate_client, web_string)

        #* CALL NLP
        test_string = ' '.join(parsed_list)
        nlp_response = self.call_nlp(nlp_client, test_string, desc_string, web_string)

        #* GENERATE KEYWORDS
        topk, topn = self.generate_keywords(nlp_response)

        #* SAVE KEYWORDS TO FIREBASE

        #!!!! THIS IS CURCIAL <CONTAINTS DEPRECATED CONTENT>
        if database == 'general':

            if path[0] == '/':
                path = path[1:]
            first_slash = path.find('/')
            ecoid = path[first_slash+1:path.find('/', first_slash+1)]

            files_ref = firebase_db.collection(u'ecosystemMatching/' + ecoid + '/keywords')
            query_ref = files_ref.where(u'docid', u'==', docid).limit(1).get()

            new_doc = {
                    u'topk': topk,
                    u'topn': topn,
                    u'docid': docid
                }

            count = 0
            for doc in query_ref:
                query_doc = doc.to_dict()
                doc.reference.update(new_doc)
                count += 1
                break

            if count == 0:
                files_ref.add(new_doc)

        #!!!! THIS IS CURCIAL <CONTAINTS DEPRECATED CONTENT>
        elif database == 'staging' or database == 'open':

            #* ADD KEYWORDS TO KEYWORDS COLLECTION
            files_ref = firebase_db.collection(u'keywords')
            query_ref = files_ref.where(u'docid', u'==', docid).limit(1).get()

            new_doc = {
                    u'topk': topk,
                    u'topn': topn,
                    u'docid': docid,
                    u'label': label,
                    u'parent': parent,
                    u'timestamp': firestore.SERVER_TIMESTAMP
                }

            # If there is an existing document then update
            count = 0
            for doc in query_ref:
                query_doc = doc.to_dict()
                doc.reference.update(new_doc)
                count += 1
                break

            # If there is no existing document then add
            if count == 0:
                firebase_db.collection(u'keywords').add(new_doc)

            #* ADD INDUSTRY TO INDUSTRY
            if industry is not None:
                files_ref = firebase_db.collection(u'industries')
                query_ref = files_ref.where(u'docid', u'==', docid).limit(1).get()
                new_doc = {
                        u'docid': docid,
                        u'industry': industry,
                        u'timestamp': firestore.SERVER_TIMESTAMP
                    }
                # If there is an existing document then update
                count = 0
                for doc in query_ref:
                    query_doc = doc.to_dict()
                    doc.reference.update(new_doc)
                    count += 1
                    break
                # If there is no existing document then add
                if count == 0:
                    firebase_db.collection(u'industries').add(new_doc)

        # print(topk)
        # print(topn)

        return "success", 200


# param = {'database':'general','docid':'Vn8DIhG9iVMckSbCDlUW','path':'/ecosystems/s2GDDevZc69Dfg8mOdOj/699c0550-dd7b-11ea-9705-1dd4eb3a9ec8/Vn8DIhG9iVMckSbCDlUW'}

# x = requests.post('http://localhost:8081/full_processing',json=param)
# x = requests.post('http://35.192.131.197/full_processing',json=param)